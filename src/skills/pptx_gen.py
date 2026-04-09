import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from src.state import BrandContext

def hex_to_rgb(hex_color):
    """Converts HEX string to RGBColor object."""
    if not hex_color:
        return RGBColor(125, 51, 255)
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(125, 51, 255) # Default purple
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def generate_pptx(content: str, brand: BrandContext, output_path: str = "output.pptx", image_assets: dict = None) -> str:
    """
    Generates a high-quality, brand-compliant PPTX presentation.
    Uses structured SLIDE_START / SLIDE_END tags.
    """
    prs = Presentation()
    
    # Brand Visuals
    primary_color = hex_to_rgb(brand.primary_color)
    secondary_color = hex_to_rgb(brand.secondary_color)
    font_family = brand.font_family or "Arial"

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Style Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = primary_color

    title = slide.shapes.title
    title.text = brand.name
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.name = font_family
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = f"Empowering Brand Communication\n{brand.tone} Tone"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 230, 230)

    # 2. Parse and Create Slides
    # Look for SLIDE_START ... SLIDE_END blocks
    slide_pattern = re.compile(r"SLIDE_START(.*?)SLIDE_END", re.DOTALL)
    slides_data = slide_pattern.findall(content)

    if not slides_data:
        # Fallback to old splitting if structure isn't found
        slides_data = content.split("\n\n")

    bullet_slide_layout = prs.slide_layouts[1]

    for data in slides_data:
        # Extract fields
        title_match = re.search(r"TITLE:\s*(.*)", data)
        content_match = re.search(r"CONTENT:\s*(.*?)(?=IMAGE_PROMPT:|$)", data, re.DOTALL | re.IGNORECASE)
        image_match = re.search(r"IMAGE_PROMPT:\s*(.*)", data, re.IGNORECASE)

        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Apply Brand Color to Title
        title_shape = slide.shapes.title
        title_text = title_match.group(1).strip() if title_match else "Key Insight"
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
        title_shape.text_frame.paragraphs[0].font.name = font_family

        # Content
        body_shape = slide.placeholders[1]
        body_text = content_match.group(1).strip() if content_match else data.strip()
        body_shape.text = body_text
        
        for p in body_shape.text_frame.paragraphs:
            p.font.size = Pt(20)
            p.font.name = font_family
            p.font.color.rgb = RGBColor(50, 50, 50)

        # Placeholder for Image (V2.1 - Will be replaced by real asset once generated)
        if image_match:
            img_prompt = image_match.group(1).strip()
            
            if image_assets and img_prompt in image_assets:
                img_path = image_assets[img_prompt]
                left = Inches(5.5)
                top = Inches(2)
                width = Inches(4)
                try:
                    slide.shapes.add_picture(img_path, left, top, width=width)
                except Exception as e:
                    print(f"Warning: Failed to insert image on slide: {e}")
            else:
                # Placeholder for Image
                left = Inches(6)
                top = Inches(2)
                width = Inches(3.5)
                txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
                tf = txBox.text_frame
                tf.text = f"[Visual Asset Recommended:\n{img_prompt[:50]}...]"
                tf.paragraphs[0].font.size = Pt(10)
                tf.paragraphs[0].font.italic = True

    # 3. Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    
    print(f"Branded PPTX generated at: {output_path}")
    return os.path.abspath(output_path)
